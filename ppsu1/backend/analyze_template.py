#!/usr/bin/env python3
"""
analyze_template.py — scans ANY designed .pptx and auto-detects slides that
are reusable "title + N repeating items" layouts (chevrons, numbered cards,
icon rows, etc). Produces a catalog.json consumable by engine.py.

A slide qualifies if it has >=3 top-level shapes that share the same shape
type + name-prefix + font size (PowerPoint's own duplication fingerprint),
each holding 1-2 non-empty text fields (ignoring tiny numeric badges like
"01", "1", "2"...). Slides without such a repeating cluster (title slides,
photo slides, freeform diagrams, unique one-off layouts) are skipped
automatically — they're not safe to reuse for arbitrary content.

Usage:
    python3 analyze_template.py template.pptx > catalog.json
"""
import hashlib
import json
import re
import sys
from collections import defaultdict
from pathlib import Path


def safe_shape_type(shape):
    """python-pptx raises NotImplementedError for shapes it cannot classify
    (ink drawings, 3D models, some smart-art internals). One exotic shape in
    an uploaded deck must not crash the whole analysis - treat it as type
    None: inert decoration that no detector matches."""
    try:
        return shape.shape_type
    except NotImplementedError:
        return None


from pptx import Presentation

BADGE_RE = re.compile(r"^\d{1,3}$")

# "SME Place Holder" boxes appear on Introduction / Learning Objectives /
# Summary designs. They are chrome the SME fills in later, NOT content, so
# they are ignored when detecting a design (and left untouched in output).
SME_RE = re.compile(r"^sme\s*place\s*holder$", re.I)


def is_sme_placeholder(shape):
    try:
        return (shape.has_text_frame
                and SME_RE.match(shape.text_frame.text.strip()) is not None)
    except Exception:
        return False

# Bump this whenever the catalog's shape changes so app.py rebuilds cached
# catalogs instead of serving a stale one (e.g. missing the new title_text).
# v3: 2-item designs are now detected (previously only 3+).
# v4: item box dimensions recorded (item_dims) for text-fit aware selection.
# v5: per-item icon slots detected (icon_ids) for icon-capable designs.
# v6: icon slots found inside item groups; SME Place Holder boxes ignored so
#     Introduction / Learning Objectives designs are detected.
# v7: section designs detected — empty banner shapes (scroll/ribbon) paired to
#     the item bodies (label_ids); reserved for sectioned raw slides.
# v8: banner shapes carrying placeholder text ("Sub heading" / "Heading")
#     recognised as label banners too, and kept out of content clustering.
# v9: section designs with the placeholder banner INSIDE each item group
#     (two-field items whose label sample is "Sub heading"); slides with
#     placeholders that could not be paired to items are skipped so the
#     sample text never leaks into a finished deck.
# v10: group-wrapped placeholder banners (banner + chip grouped together)
#      recognised as label banners.
# v11: content_top recorded — how far down a design's artwork starts, so a raw
#      slide carrying an intro box above its points can keep that box.
# v12: content_top now measures the item boxes too (v11 wrongly excluded them,
#      so designs whose labels sit high were judged clear).
# v13: shapes holding ONLY a number badge ("01") are decoration and no longer
#      join the item clustering, so designs that draw the badge separately
#      from its text are detected.
CATALOG_VERSION = 13


def content_top(slide, slide_area, skip_ids=()):
    """Top edge (EMU) of the highest piece of a design's artwork, ignoring the
    title/subtitle and any near-full-canvas background. Everything below this
    line is the design; anything a raw slide wants to keep above it (an intro
    box) therefore has room."""
    tops = []
    for s in slide.shapes:
        if s.shape_id in skip_ids:
            continue
        w, h = int(s.width or 0), int(s.height or 0)
        if w <= 0 or h <= 0:
            continue
        if w * h > 0.55 * slide_area:      # background / full-bleed artwork
            continue
        tops.append(int(s.top or 0))
    return min(tops) if tops else 0


def slide_has_label_placeholder(slide):
    """True when ANY shape on the slide (recursing into groups) carries
    placeholder label text ("Sub heading" / "Heading")."""
    def walk(s):
        if is_label_placeholder(s):
            return True
        if safe_shape_type(s) == 6:
            return any(walk(c) for c in s.shapes)
        return False
    return any(walk(s) for s in slide.shapes)

# sample text an author types into a banner to mark where the section
# subheading goes — placeholder, not content
LABEL_PLACEHOLDER_RE = re.compile(r"^(sub\s*-?\s*heading|heading|label)$",
                                  re.I)


def is_label_placeholder(shape):
    """The shape is a placeholder label banner: its own text is "Sub
    heading"/"Heading", or it is a GROUP whose only real text (ignoring
    number badges) is such a placeholder — authors often group the banner
    with a chip or oval."""
    try:
        if (shape.has_text_frame and LABEL_PLACEHOLDER_RE.match(
                shape.text_frame.text.strip()) is not None):
            return True
        if safe_shape_type(shape) == 6:
            leaves = collect_text_leaves(shape)
            return (len(leaves) == 1
                    and LABEL_PLACEHOLDER_RE.match(leaves[0][1]) is not None)
    except Exception:
        pass
    return False


def _is_label_banner(s):
    """A banner shape that receives a section subheading: an autoshape whose
    whole text is a placeholder ("Sub heading"), or a text-less autoshape
    drawn as a scroll/ribbon or explicitly named for the job."""
    if is_label_placeholder(s):
        return True
    if shape_has_text(s):
        return False
    nm = name_prefix(s.name).lower()
    if nm.split(":")[0].strip() in ("scroll", "banner", "label", "heading"):
        return True
    for el in s._element.iter():
        tag = getattr(el, "tag", "")
        if isinstance(tag, str) and tag.endswith("prstGeom"):
            prst = (el.get("prst") or "").lower()
            if "scroll" in prst or "ribbon" in prst:
                return True
    return False


def find_label_banners(slide, ordered_items):
    """Per-item empty banner shapes (section designs): each item body gets its
    nearest unused banner. Returns shape-ids in item order, or None unless
    EVERY item has one."""
    cands = [s for s in slide.shapes if _is_label_banner(s)]
    if len(cands) < len(ordered_items):
        return None

    def center(s):
        return ((s.left or 0) + (s.width or 0) / 2,
                (s.top or 0) + (s.height or 0) / 2)

    used, out = set(), []
    for it in ordered_items:
        ix, iy = center(it)
        best, best_d = None, None
        for c in cands:
            if c.shape_id in used:
                continue
            cx, cy = center(c)
            d = (cx - ix) ** 2 + (cy - iy) ** 2
            if best_d is None or d < best_d:
                best, best_d = c, d
        if best is None:
            return None
        used.add(best.shape_id)
        out.append(best.shape_id)
    return out

# an icon slot in a hand-made template is usually an actual small picture;
# anything bigger than this (EMU, 1.2 in) is real imagery — a photo, an
# illustration, an SME picture — and must NEVER be replaced by an icon
MAX_ICON_SLOT_EMU = 1097280
MIN_ICON_SLOT_EMU = 91440          # smaller than 0.1 in is a stray artifact


def _is_icon_picture(s):
    return (safe_shape_type(s) == 13
            and MIN_ICON_SLOT_EMU <= int(s.width or 0) <= MAX_ICON_SLOT_EMU
            and MIN_ICON_SLOT_EMU <= int(s.height or 0) <= MAX_ICON_SLOT_EMU)


def _item_icon_picture(item):
    """The icon PICTURE inside an item shape (recursing into groups): the
    largest picture that is still icon-sized. None if the item has none.
    Big pictures inside the item are real images and are never returned."""
    found = []
    def walk(s):
        if _is_icon_picture(s):
            found.append(s)
        if safe_shape_type(s) == 6:
            for c in s.shapes:
                walk(c)
    walk(item)
    if not found:
        return None
    return max(found, key=lambda s: int(s.width or 0) * int(s.height or 0))


def find_icon_slots(slide, ordered_items):
    """Per-item icon slots, two detection styles:

    1. Real designed decks: each item (usually a group) carries its own small
       icon PICTURE inside it — used when EVERY item has one. Only that small
       picture is ever replaced; full-size images never qualify.
    2. Generated templates: text-less top-level shapes NAMED "Icon n" (or
       free-standing small pictures), paired to items by proximity.

    Returns the slot shape-ids in item order, or None unless EVERY item gets
    its own slot."""
    inside = [_item_icon_picture(it) for it in ordered_items]
    if all(p is not None for p in inside):
        return [p.shape_id for p in inside]

    cands = []
    for s in slide.shapes:
        if shape_has_text(s):
            continue
        named = name_prefix(s.name).lower() == "icon"
        if named or _is_icon_picture(s):
            cands.append(s)
    if len(cands) < len(ordered_items):
        return None

    def center(s):
        return ((s.left or 0) + (s.width or 0) / 2,
                (s.top or 0) + (s.height or 0) / 2)

    used, slots = set(), []
    for it in ordered_items:
        ix, iy = center(it)
        best, best_d = None, None
        for c in cands:
            if c.shape_id in used:
                continue
            cx, cy = center(c)
            d = (cx - ix) ** 2 + (cy - iy) ** 2
            if best_d is None or d < best_d:
                best, best_d = c, d
        if best is None:
            return None
        used.add(best.shape_id)
        slots.append(best.shape_id)
    return slots


def design_key(item_shapes):
    """Fingerprint the visual style of a layout's repeating items — preset
    geometry, colors (fill + line + text), and font size — while ignoring how
    many items there are and what text they hold.

    Two layouts that look the same but carry a different number of points (e.g.
    a 4-item and a 5-item version of the same chevron design, where the second
    just adds one more chevron at the bottom) therefore produce the SAME key,
    so the engine can treat them as one design and not repeat that look.
    """
    geoms, colors, fonts = set(), set(), set()
    for s in item_shapes:
        for c in s._element.iter():
            tag = c.tag
            if not isinstance(tag, str):  # skip comments / PIs
                continue
            local = tag.rsplit("}", 1)[-1]
            if local == "prstGeom" and c.get("prst"):
                geoms.add(c.get("prst"))
            elif local in ("srgbClr", "schemeClr") and c.get("val"):
                colors.add(local + ":" + c.get("val"))
        fonts.add(round(max_font_size(s)))
    payload = json.dumps({"g": sorted(geoms), "c": sorted(colors),
                          "f": sorted(fonts)}, sort_keys=True)
    return hashlib.sha1(payload.encode()).hexdigest()[:12]


def max_font_size(shape):
    best = 0.0
    def walk(s):
        nonlocal best
        if s.has_text_frame:
            for p in s.text_frame.paragraphs:
                for r in p.runs:
                    if r.font.size:
                        best = max(best, r.font.size.pt)
        if safe_shape_type(s) == 6:
            for c in s.shapes:
                walk(c)
    walk(shape)
    return best


def name_prefix(name):
    # Google Slides exports name shapes "Google Shape;<id>;p<page>"; the id is
    # per-shape noise, so collapse them all into one name family
    if (name or "").startswith("Google Shape;"):
        return "Google Shape"
    return re.sub(r"\s*\d+\s*$", "", name).strip()


def shape_has_text(shape):
    if shape.has_text_frame and shape.text_frame.text.strip():
        return True
    if safe_shape_type(shape) == 6:
        return any(shape_has_text(c) for c in shape.shapes)
    return False


def collect_text_leaves(shape):
    """Non-empty, non-badge text-bearing leaves inside `shape` (recursive)."""
    leaves = []
    def walk(s):
        if s.has_text_frame and s.text_frame.text.strip():
            txt = s.text_frame.text.strip()
            if BADGE_RE.match(txt):
                return
            leaves.append((s.shape_id, txt, max_font_size(s),
                           int(s.width or 0), int(s.height or 0)))
            return  # don't recurse into a shape that already has text
        if safe_shape_type(s) == 6:
            for c in s.shapes:
                walk(c)
    if safe_shape_type(shape) == 6:
        for c in shape.shapes:
            walk(c)
    else:
        walk(shape)
    return leaves


def find_badge_number(shape):
    def walk(s):
        if s.has_text_frame:
            t = s.text_frame.text.strip()
            if BADGE_RE.match(t):
                return int(t)
        if safe_shape_type(s) == 6:
            for c in s.shapes:
                r = walk(c)
                if r is not None:
                    return r
        return None
    return walk(shape)


def analyze_slide(slide, slide_num, slide_area):
    top = list(slide.shapes)
    # A shape whose only text is a number badge ("01") is decoration, not
    # content: some designs draw the badge as its own shape beside the point
    # text. Excluding them here lets the real point boxes cluster together.
    text_top = [s for s in top
                if shape_has_text(s) and not is_sme_placeholder(s)
                and not is_label_placeholder(s)
                and collect_text_leaves(s)]
    if len(text_top) < 2:  # need at least a title + something
        return None

    clusters = defaultdict(list)
    for s in text_top:
        key = (safe_shape_type(s), name_prefix(s.name), round(max_font_size(s)))
        clusters[key].append(s)
    # >= 2 so a genuine 2-item design (e.g. a two-box "left / right" layout) is
    # detected too, not just 3+ item designs. A bare pair is weaker evidence, so
    # it is guarded further below (it must have a clearly larger title above it).
    candidates = [v for v in clusters.values() if len(v) >= 2]
    # drop clusters that are pure decoration (e.g. number badges only) once
    # badge text is excluded — they have nothing left to hold raw content
    candidates = [v for v in candidates if all(collect_text_leaves(s) for s in v)]

    if not candidates:
        # fallback: title + exactly one body shape (e.g. a single content
        # block/card) — usable as a 1-item template that absorbs all bullets
        remaining = [s for s in text_top if s is not None]
        remaining.sort(key=lambda s: (-max_font_size(s), s.top or 0))
        if len(remaining) == 2 and collect_text_leaves(remaining[1]):
            title_shape, body_shape = remaining
            # a single body box that itself already holds several bullet lines
            # is NOT a safe 1-item template: reusing it forces all of a raw
            # slide's points to be merged into one paragraph. Skip such slides.
            if body_shape.has_text_frame:
                body_paras = [p for p in body_shape.text_frame.paragraphs
                              if p.text.strip()]
                if len(body_paras) >= 2:
                    return None
            body_area = (body_shape.width or 0) * (body_shape.height or 0)
            if body_area / slide_area < 0.06:  # too small to be a real content block (e.g. a name tag)
                return None
            if slide_has_label_placeholder(slide):
                return None
            return {
                "slide_file": f"slide{slide_num}.xml",
                "title_id": title_shape.shape_id,
                "title_text": title_shape.text_frame.text.strip(),
                "subtitle_id": None,
                "item_count": 1,
                "item_fields": 1,
                "items": [body_shape.shape_id],
                "item_dims": [[int(body_shape.width or 0),
                               int(body_shape.height or 0)]],
                "content_top": content_top(slide, slide_area,
                                           {title_shape.shape_id}),
                "design_key": design_key([body_shape]),
            }
        return None
    items = max(candidates, key=len)

    # if most items carry a numeric badge but one or two don't, those are
    # probably a stray label box (e.g. a subtitle) that happened to share
    # the same shape signature — not a genuine repeating item
    badge_presence = [find_badge_number(s) is not None for s in items]
    if any(badge_presence) and not all(badge_presence):
        items = [s for s, has in zip(items, badge_presence) if has]
        if len(items) < 3:
            return None

    remaining = [s for s in text_top if s not in items]
    if not remaining:
        return None
    remaining.sort(key=lambda s: (-max_font_size(s), s.top or 0))
    title_shape = remaining[0]

    item_font = round(max_font_size(items[0]))

    # a PAIR of items is only trusted when a clearly larger title sits above it.
    # This stops a title slide's two similarly-styled lines (or a title+subtitle)
    # from being mistaken for a 2-item content design. 3+ items are unambiguous
    # on their own and skip this check, so existing templates are unaffected.
    if len(items) == 2 and round(max_font_size(title_shape)) <= item_font:
        return None
    subtitle_shape = None
    subtitle_pool = [s for s in remaining[1:] if max_font_size(s) > item_font]
    if subtitle_pool:
        subtitle_pool.sort(key=lambda s: s.top or 0)
        subtitle_shape = subtitle_pool[0]

    positions = [(s.top, s.left) for s in items]
    if len(set(positions)) == len(items):
        ordered = sorted(items, key=lambda s: (s.top, s.left))
    else:
        badges = [(find_badge_number(s), s) for s in items]
        badge_vals = {b for b, _ in badges}
        if None not in badge_vals and len(badge_vals) == len(items):
            ordered = [s for _, s in sorted(badges, key=lambda x: x[0])]
        else:
            # can't reliably order these items — too risky to use as a template
            return None

    # every top-level text-bearing shape must be accounted for as title,
    # subtitle, or an item — an extra unclassified shape means original
    # template copy would leak into the output untouched
    accounted_ids = {s.shape_id for s in items} | {title_shape.shape_id}
    if subtitle_shape is not None:
        accounted_ids.add(subtitle_shape.shape_id)
    stray = [s for s in text_top if s.shape_id not in accounted_ids]
    # a stray shape is only a problem if it holds real content (not just a
    # decorative number badge) — that content would leak into the output untouched
    if any(collect_text_leaves(s) for s in stray):
        return None

    def _is_ph_leaf(leaf):
        return LABEL_PLACEHOLDER_RE.match(leaf[1]) is not None

    item_fields = []
    for s in ordered:
        leaves = collect_text_leaves(s)
        if not leaves:
            return None
        leaves.sort(key=lambda x: -x[2])
        pair = leaves[:2]
        # a "Sub heading" placeholder leaf inside the item is ALWAYS the
        # label field, whatever its font size says
        if len(pair) == 2 and _is_ph_leaf(pair[1]) and not _is_ph_leaf(pair[0]):
            pair = [pair[1], pair[0]]
        item_fields.append(pair)

    nfields = {len(f) for f in item_fields}
    if len(nfields) != 1:
        return None
    nfields = nfields.pop()

    # section design, two authoring styles:
    # * two-field items whose label sample is a placeholder ("Sub heading"
    #   drawn INSIDE each item group);
    # * single-field bodies with a separate banner shape per item.
    two_field_section = (nfields == 2
                         and all(_is_ph_leaf(f[0]) for f in item_fields))
    label_slots = (find_label_banners(slide, ordered)
                   if nfields == 1 else None)
    is_section = two_field_section or bool(label_slots)

    # a slide carrying placeholder banners that could NOT be paired to its
    # items (e.g. one banner for two bullet rows) is unusable — reusing it
    # would leave literal "Sub heading" text in the finished deck
    if not is_section and slide_has_label_placeholder(slide):
        return None

    return {
        "slide_file": f"slide{slide_num}.xml",
        "title_id": title_shape.shape_id,
        "title_text": title_shape.text_frame.text.strip(),
        "subtitle_id": subtitle_shape.shape_id if subtitle_shape else None,
        "item_count": len(items),
        "item_fields": nfields,
        "items": [f[0][0] if nfields == 1 else [f[0][0], f[1][0]] for f in item_fields],
        # width/height of the box that receives the point text (the description
        # box on two-field designs) — lets the engine judge whether a given text
        # actually FITS this design before choosing it
        "item_dims": [[f[0][3], f[0][4]] if nfields == 1 else [f[1][3], f[1][4]]
                      for f in item_fields],
        # one icon slot per item (shape ids, item order) when the design has
        # them; None for text-only designs. Icon designs get a distinct
        # design_key suffix so rotation never confuses an icon design with
        # its text-only look-alike.
        "icon_ids": (icon_slots := find_icon_slots(slide, ordered)),
        # one banner per item (single-field section designs): the banner
        # receives a section's bold subheading, the item body its bullets.
        # Two-field section designs carry the label INSIDE each item instead.
        "label_ids": label_slots,
        "section_design": is_section,
        # where this design's artwork starts — used to decide whether a raw
        # slide's intro box can be kept above it. Only the title/subtitle are
        # excluded: item boxes DO count, since a design may place its labels
        # above its artwork.
        "content_top": content_top(
            slide, slide_area,
            {title_shape.shape_id} | (
                {subtitle_shape.shape_id} if subtitle_shape else set())),
        "design_key": (design_key(ordered)
                       + ("-icons" if icon_slots else "")
                       + ("-sections" if is_section else "")),
    }


def build_catalog(template_path):
    prs = Presentation(str(template_path))
    slide_area = prs.slide_width * prs.slide_height
    patterns = defaultdict(list)
    skipped = []
    for i, slide in enumerate(prs.slides, 1):
        entry = analyze_slide(slide, i, slide_area)
        if entry:
            patterns[entry["item_count"]].append(entry)
        else:
            skipped.append(i)
    return {
        "version": CATALOG_VERSION,
        "source_deck": Path(template_path).name,
        "patterns": {str(k): v for k, v in sorted(patterns.items())},
        "skipped_slides": skipped,
    }


if __name__ == "__main__":
    if len(sys.argv) != 2:
        print("Usage: analyze_template.py template.pptx", file=sys.stderr)
        sys.exit(1)
    catalog = build_catalog(sys.argv[1])
    print(json.dumps(catalog, indent=2))
