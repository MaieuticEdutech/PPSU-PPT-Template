#!/usr/bin/env python3
"""
analyze_template.py — scans ANY designed .pptx and auto-detects slides that
are reusable "title + N repeating items" layouts (chevrons, numbered cards,
icon rows, etc). Produces a catalog.json consumable by engine.py.

A slide qualifies if it has >=3 top-level shapes that share the same shape
type + name-prefix + font size (PowerPoint's own duplication fingerprint),
each holding 1-2 non-empty text fields (ignoring tiny numeric badges like
"01", "1", "2"...). Slides without such a repeating cluster (title slides,
photo slides, freeform diagrams, unique one-off layouts) are skipped
automatically — they're not safe to reuse for arbitrary content.

Usage:
    python3 analyze_template.py template.pptx > catalog.json
"""
import hashlib
import json
import re
import sys
from collections import defaultdict
from pathlib import Path


def safe_shape_type(shape):
    """python-pptx raises NotImplementedError for shapes it cannot classify
    (ink drawings, 3D models, some smart-art internals). One exotic shape in
    an uploaded deck must not crash the whole analysis - treat it as type
    None: inert decoration that no detector matches."""
    try:
        return shape.shape_type
    except NotImplementedError:
        return None


from pptx import Presentation

BADGE_RE = re.compile(r"^\d{1,3}$")

# Bump this whenever the catalog's shape changes so app.py rebuilds cached
# catalogs instead of serving a stale one (e.g. missing the new title_text).
# v3: 2-item designs are now detected (previously only 3+).
# v4: item box dimensions recorded (item_dims) for text-fit aware selection.
CATALOG_VERSION = 4


def design_key(item_shapes):
    """Fingerprint the visual style of a layout's repeating items — preset
    geometry, colors (fill + line + text), and font size — while ignoring how
    many items there are and what text they hold.

    Two layouts that look the same but carry a different number of points (e.g.
    a 4-item and a 5-item version of the same chevron design, where the second
    just adds one more chevron at the bottom) therefore produce the SAME key,
    so the engine can treat them as one design and not repeat that look.
    """
    geoms, colors, fonts = set(), set(), set()
    for s in item_shapes:
        for c in s._element.iter():
            tag = c.tag
            if not isinstance(tag, str):  # skip comments / PIs
                continue
            local = tag.rsplit("}", 1)[-1]
            if local == "prstGeom" and c.get("prst"):
                geoms.add(c.get("prst"))
            elif local in ("srgbClr", "schemeClr") and c.get("val"):
                colors.add(local + ":" + c.get("val"))
        fonts.add(round(max_font_size(s)))
    payload = json.dumps({"g": sorted(geoms), "c": sorted(colors),
                          "f": sorted(fonts)}, sort_keys=True)
    return hashlib.sha1(payload.encode()).hexdigest()[:12]


def max_font_size(shape):
    best = 0.0
    def walk(s):
        nonlocal best
        if s.has_text_frame:
            for p in s.text_frame.paragraphs:
                for r in p.runs:
                    if r.font.size:
                        best = max(best, r.font.size.pt)
        if safe_shape_type(s) == 6:
            for c in s.shapes:
                walk(c)
    walk(shape)
    return best


def name_prefix(name):
    return re.sub(r"\s*\d+\s*$", "", name).strip()


def shape_has_text(shape):
    if shape.has_text_frame and shape.text_frame.text.strip():
        return True
    if safe_shape_type(shape) == 6:
        return any(shape_has_text(c) for c in shape.shapes)
    return False


def collect_text_leaves(shape):
    """Non-empty, non-badge text-bearing leaves inside `shape` (recursive)."""
    leaves = []
    def walk(s):
        if s.has_text_frame and s.text_frame.text.strip():
            txt = s.text_frame.text.strip()
            if BADGE_RE.match(txt):
                return
            leaves.append((s.shape_id, txt, max_font_size(s),
                           int(s.width or 0), int(s.height or 0)))
            return  # don't recurse into a shape that already has text
        if safe_shape_type(s) == 6:
            for c in s.shapes:
                walk(c)
    if safe_shape_type(shape) == 6:
        for c in shape.shapes:
            walk(c)
    else:
        walk(shape)
    return leaves


def find_badge_number(shape):
    def walk(s):
        if s.has_text_frame:
            t = s.text_frame.text.strip()
            if BADGE_RE.match(t):
                return int(t)
        if safe_shape_type(s) == 6:
            for c in s.shapes:
                r = walk(c)
                if r is not None:
                    return r
        return None
    return walk(shape)


def analyze_slide(slide, slide_num, slide_area):
    top = list(slide.shapes)
    text_top = [s for s in top if shape_has_text(s)]
    if len(text_top) < 2:  # need at least a title + something
        return None

    clusters = defaultdict(list)
    for s in text_top:
        key = (safe_shape_type(s), name_prefix(s.name), round(max_font_size(s)))
        clusters[key].append(s)
    # >= 2 so a genuine 2-item design (e.g. a two-box "left / right" layout) is
    # detected too, not just 3+ item designs. A bare pair is weaker evidence, so
    # it is guarded further below (it must have a clearly larger title above it).
    candidates = [v for v in clusters.values() if len(v) >= 2]
    # drop clusters that are pure decoration (e.g. number badges only) once
    # badge text is excluded — they have nothing left to hold raw content
    candidates = [v for v in candidates if all(collect_text_leaves(s) for s in v)]

    if not candidates:
        # fallback: title + exactly one body shape (e.g. a single content
        # block/card) — usable as a 1-item template that absorbs all bullets
        remaining = [s for s in text_top if s is not None]
        remaining.sort(key=lambda s: (-max_font_size(s), s.top or 0))
        if len(remaining) == 2 and collect_text_leaves(remaining[1]):
            title_shape, body_shape = remaining
            # a single body box that itself already holds several bullet lines
            # is NOT a safe 1-item template: reusing it forces all of a raw
            # slide's points to be merged into one paragraph. Skip such slides.
            if body_shape.has_text_frame:
                body_paras = [p for p in body_shape.text_frame.paragraphs
                              if p.text.strip()]
                if len(body_paras) >= 2:
                    return None
            body_area = (body_shape.width or 0) * (body_shape.height or 0)
            if body_area / slide_area < 0.06:  # too small to be a real content block (e.g. a name tag)
                return None
            return {
                "slide_file": f"slide{slide_num}.xml",
                "title_id": title_shape.shape_id,
                "title_text": title_shape.text_frame.text.strip(),
                "subtitle_id": None,
                "item_count": 1,
                "item_fields": 1,
                "items": [body_shape.shape_id],
                "item_dims": [[int(body_shape.width or 0),
                               int(body_shape.height or 0)]],
                "design_key": design_key([body_shape]),
            }
        return None
    items = max(candidates, key=len)

    # if most items carry a numeric badge but one or two don't, those are
    # probably a stray label box (e.g. a subtitle) that happened to share
    # the same shape signature — not a genuine repeating item
    badge_presence = [find_badge_number(s) is not None for s in items]
    if any(badge_presence) and not all(badge_presence):
        items = [s for s, has in zip(items, badge_presence) if has]
        if len(items) < 3:
            return None

    remaining = [s for s in text_top if s not in items]
    if not remaining:
        return None
    remaining.sort(key=lambda s: (-max_font_size(s), s.top or 0))
    title_shape = remaining[0]

    item_font = round(max_font_size(items[0]))

    # a PAIR of items is only trusted when a clearly larger title sits above it.
    # This stops a title slide's two similarly-styled lines (or a title+subtitle)
    # from being mistaken for a 2-item content design. 3+ items are unambiguous
    # on their own and skip this check, so existing templates are unaffected.
    if len(items) == 2 and round(max_font_size(title_shape)) <= item_font:
        return None
    subtitle_shape = None
    subtitle_pool = [s for s in remaining[1:] if max_font_size(s) > item_font]
    if subtitle_pool:
        subtitle_pool.sort(key=lambda s: s.top or 0)
        subtitle_shape = subtitle_pool[0]

    positions = [(s.top, s.left) for s in items]
    if len(set(positions)) == len(items):
        ordered = sorted(items, key=lambda s: (s.top, s.left))
    else:
        badges = [(find_badge_number(s), s) for s in items]
        badge_vals = {b for b, _ in badges}
        if None not in badge_vals and len(badge_vals) == len(items):
            ordered = [s for _, s in sorted(badges, key=lambda x: x[0])]
        else:
            # can't reliably order these items — too risky to use as a template
            return None

    # every top-level text-bearing shape must be accounted for as title,
    # subtitle, or an item — an extra unclassified shape means original
    # template copy would leak into the output untouched
    accounted_ids = {s.shape_id for s in items} | {title_shape.shape_id}
    if subtitle_shape is not None:
        accounted_ids.add(subtitle_shape.shape_id)
    stray = [s for s in text_top if s.shape_id not in accounted_ids]
    # a stray shape is only a problem if it holds real content (not just a
    # decorative number badge) — that content would leak into the output untouched
    if any(collect_text_leaves(s) for s in stray):
        return None

    item_fields = []
    for s in ordered:
        leaves = collect_text_leaves(s)
        if not leaves:
            return None
        leaves.sort(key=lambda x: -x[2])
        item_fields.append(leaves[:2])

    nfields = {len(f) for f in item_fields}
    if len(nfields) != 1:
        return None
    nfields = nfields.pop()

    return {
        "slide_file": f"slide{slide_num}.xml",
        "title_id": title_shape.shape_id,
        "title_text": title_shape.text_frame.text.strip(),
        "subtitle_id": subtitle_shape.shape_id if subtitle_shape else None,
        "item_count": len(items),
        "item_fields": nfields,
        "items": [f[0][0] if nfields == 1 else [f[0][0], f[1][0]] for f in item_fields],
        # width/height of the box that receives the point text (the description
        # box on two-field designs) — lets the engine judge whether a given text
        # actually FITS this design before choosing it
        "item_dims": [[f[0][3], f[0][4]] if nfields == 1 else [f[1][3], f[1][4]]
                      for f in item_fields],
        "design_key": design_key(ordered),
    }


def build_catalog(template_path):
    prs = Presentation(str(template_path))
    slide_area = prs.slide_width * prs.slide_height
    patterns = defaultdict(list)
    skipped = []
    for i, slide in enumerate(prs.slides, 1):
        entry = analyze_slide(slide, i, slide_area)
        if entry:
            patterns[entry["item_count"]].append(entry)
        else:
            skipped.append(i)
    return {
        "version": CATALOG_VERSION,
        "source_deck": Path(template_path).name,
        "patterns": {str(k): v for k, v in sorted(patterns.items())},
        "skipped_slides": skipped,
    }


if __name__ == "__main__":
    if len(sys.argv) != 2:
        print("Usage: analyze_template.py template.pptx", file=sys.stderr)
        sys.exit(1)
    catalog = build_catalog(sys.argv[1])
    print(json.dumps(catalog, indent=2))
