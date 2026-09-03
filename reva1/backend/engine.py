#!/usr/bin/env python3
"""
ppt-template-engine: converts a raw content PPT into a finished PPT by
reusing an existing designed PPT's slide layouts and swapping in only the
text — no other design work.

Usage:
    python3 engine.py --raw raw.pptx --template template.pptx --catalog catalog.json --out final.pptx

Self-contained: only depends on python-pptx + the standard library (no
sandbox-only helper scripts). The catalog maps "number of bullets" -> which
template slide + which shapes in it hold the title / subtitle / per-item
text. It is built once per template deck (see analyze_template.py) and
reused for any raw deck matched against that same template.
"""
import argparse
import copy
import json
import random
import re
import sys
from collections import Counter
from pathlib import Path


def safe_shape_type(shape):
    """python-pptx raises NotImplementedError for shapes it cannot classify
    (ink drawings, 3D models, some smart-art internals). One exotic shape in
    an uploaded deck must not crash the whole analysis - treat it as type
    None: inert decoration that no detector matches."""
    try:
        return shape.shape_type
    except NotImplementedError:
        return None


from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml.ns import qn
from pptx.parts.slide import SlidePart
from pptx.util import Emu, Pt

R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


RAW_BADGE_RE = re.compile(r"^\d{1,3}$")

# MSO_SHAPE_TYPE.TEXT_BOX - a plain text box, as opposed to an auto-shape
# (rounded box / callout) that carries text as part of a diagram.
TEXT_BOX_TYPE = 17

# Two shapes whose tops are within this distance (EMU; 914400 = 1 inch) sit on
# the same visual row, i.e. side by side rather than stacked in a list.
SAME_ROW_EMU = 100000

# How far the points' LEFT edges may differ and still count as one aligned
# column (EMU; 457200 = 0.5 inch). This is the decisive test: a list shares a
# left edge, while a diagram's boxes sit 4+ inches apart horizontally.
ALIGN_TOL_EMU = 457200

# Widths are allowed to vary more (914400 = 1 inch): authors routinely size
# individual point boxes a little differently, and width alone never tells a
# list apart from a diagram — the left-edge test above does that.
WIDTH_TOL_EMU = 914400

# "Section" slides that must ALWAYS use their own dedicated template design
# (matched by the design slide's title text) rather than a rotated/random
# layout. A raw slide whose heading is one of these is pinned to the template
# design carrying the same title; and those dedicated designs are withheld from
# the general rotation so ordinary content never lands on them.
RESERVED_TITLES = {
    "summary",
    "learning objectives", "learning objective", "objectives",
}


def _norm_title(text):
    """Normalise a heading/title for matching: lowercase, single-spaced."""
    return re.sub(r"\s+", " ", (text or "").strip().lower())


def _para_texts(shape):
    """Non-empty paragraph texts of a shape."""
    return [p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()]


def _shape_font(shape):
    """Largest run font size (pt) in a shape, or 0.0 if none is set."""
    best = 0.0
    if shape.has_text_frame:
        for p in shape.text_frame.paragraphs:
            for r in p.runs:
                if r.font.size:
                    best = max(best, r.font.size.pt)
    return best


def _extract_multibox(content):
    """Type B raw slide: heading + subheading + one text box per bullet.

    `content` is the slide's text-bearing shapes with decorative number badges
    already removed, sorted by (top, left). The topmost box is the heading; the
    bullets are the largest group of boxes that share the same font size (the
    repeating point boxes); a single larger-font box sitting above that group is
    treated as the subheading.
    """
    heading = content[0].text_frame.text.strip()
    rest = content[1:]
    if not rest:
        return heading, None, []

    fonts = [(s, round(_shape_font(s))) for s in rest]
    counts = Counter(f for _, f in fonts)
    top_count = max(counts.values())
    # the bullet font is the most common size among the remaining boxes; on a
    # tie prefer the smaller size (point text is smaller and more numerous than
    # a lone subheading)
    bullet_font = min(f for f, c in counts.items() if c == top_count)
    bullet_shapes = [s for s, f in fonts if f == bullet_font]
    others = [s for s, f in fonts if f != bullet_font]

    subheading = None
    if others:
        top_bullet = min((s.top or 0) for s in bullet_shapes)
        above = [s for s in others if (s.top or 0) <= top_bullet]
        if above:
            above.sort(key=lambda s: -_shape_font(s))
            subheading = above[0].text_frame.text.strip()

    bullet_shapes.sort(key=lambda s: (s.top or 0, s.left or 0))
    bullets = [s.text_frame.text.strip() for s in bullet_shapes]
    return heading, subheading, bullets


def _extract_named_items(item_shapes, others):
    """Type C raw slide: a cluster of identically-named point boxes (e.g. several
    shapes all named "PointText"), each holding an optional bold label followed
    by a description. The non-item text boxes give the heading (topmost) and,
    if one sits above the items, the subheading.

    Returns (heading, subheading, bullets, pairs) where bullets[i] is the item's
    combined "label: description" text (so a single-field design still shows
    both) and pairs[i] is (label, description) — or None when the item has no
    separate label — so a two-field design can place them in distinct boxes.
    """
    others = sorted(others, key=lambda s: (s.top or 0, s.left or 0))
    heading = others[0].text_frame.text.strip() if others else ""

    item_shapes = sorted(item_shapes, key=lambda s: (s.top or 0, s.left or 0))
    item_top = min((s.top or 0) for s in item_shapes)

    subheading = None
    if len(others) >= 2 and (others[1].top or 0) < item_top:
        subheading = others[1].text_frame.text.strip()

    bullets, pairs = [], []
    for s in item_shapes:
        paras = _para_texts(s)
        if not paras:
            continue
        if len(paras) >= 2:
            label = paras[0].rstrip(":").strip()
            desc = " ".join(paras[1:])
            bullets.append(f"{label}: {desc}")
            pairs.append((label, desc))
        else:
            bullets.append(paras[0])
            pairs.append(None)
    return heading, subheading, bullets, pairs


# Some decks name their point boxes explicitly. That name is an authored
# signal of intent, so a cluster of just TWO is trustworthy there. Any other
# repeated name still needs 3+, because exactly two same-named shapes are
# usually a title slide's two lines (both named "Text 2" in these decks).
EXPLICIT_POINT_NAME = "PointText"


def _is_point_column(items, min_items=3):
    """True when `items` really are a plain vertical list of point boxes.

    Guards against a repeated shape name (or font size) that actually belongs
    to a diagram: e.g. eight auto-shapes all named "DiagBox" holding ClassA /
    ClassB / ClassC laid out side by side. A genuine list is plain text boxes,
    stacked one under another, sharing a left edge and a width.
    """
    if len(items) < min_items:
        return False
    # NB: shape type is deliberately NOT checked here - some decks build a
    # perfectly ordinary list out of auto-shapes ("CalloutBox"). Geometry alone
    # separates a list from a diagram, and does so decisively.
    # stacked vertically - two boxes on one row means a side-by-side layout
    tops = sorted(s.top or 0 for s in items)
    if any(tops[i + 1] - tops[i] < SAME_ROW_EMU for i in range(len(tops) - 1)):
        return False
    # one aligned column - diagram labels scatter horizontally at odd widths
    lefts = [s.left or 0 for s in items]
    widths = [s.width or 0 for s in items]
    if max(lefts) - min(lefts) > ALIGN_TOL_EMU:
        return False
    if max(widths) - min(widths) > WIDTH_TOL_EMU:
        return False
    return True


def _single_body_box(content, slide_shapes, decor):
    """Type E raw slide: a title (+ optional subtitle) and ONE text box that
    holds every point as its own paragraph.

    Used by decks that keep all the bullets in a single box under generic shape
    names ("TextBox 3"), so neither the named-box detector (Type A) nor the
    repeated-box detectors (Types B/C/D) apply. Guarded the same way as those:
    all text must sit in plain text boxes, every other text shape must be
    explainable as the title/subtitle above or a small footer below, and any
    real graphic means the slide is left exactly as it is.

    Returns (body_shape, other_shapes) or None.
    """
    if len(content) < 2:
        return None
    # text inside an auto-shape/group is a diagram or callout, not a point list
    if any(safe_shape_type(s) != TEXT_BOX_TYPE for s in content):
        return None

    multi = [s for s in content if len(_para_texts(s)) >= 2]
    if len(multi) != 1:  # exactly one box may hold the points
        return None
    body = multi[0]
    others = [s for s in content if s is not body]
    if not others:  # nothing left to act as the heading
        return None

    body_font = round(_shape_font(body))
    body_top = body.top or 0
    body_bottom = body_top + (body.height or 0)
    for s in others:
        f = round(_shape_font(s))
        top = s.top or 0
        # title/subtitle above the points: larger font, OR a single-line
        # box sharing the body font (some decks set subheading and bullets
        # at the same size — same font, but one line above the box)
        if top <= body_top and (f > body_font
                                 or (f == body_font
                                     and len(_para_texts(s)) == 1)):
            continue
        if top >= body_bottom and f < body_font:
            continue  # small footer underneath
        return None   # unexplained text -> keep the slide as it is

    # a graphic that is neither deck chrome nor sitting alongside the body text
    # is real content (a diagram, a chart) that a rebuild would destroy
    for s in slide_shapes:
        if s.has_text_frame and s.text_frame.text.strip():
            continue
        if _shape_sig(s) in decor:
            continue
        top = s.top or 0
        if top < body_bottom and top + (s.height or 0) > body_top:
            continue
        return None
    return body, others


def _skip_reason(slide, content):
    """Plain-English explanation of why a slide was NOT turned into a design.

    Read-only: it inspects the slide and reports the first condition that
    explains the skip, in the same order the detectors apply their guards. It
    never influences the decision itself, so it cannot change what gets
    redesigned.
    """
    if not content:
        return "no text content — image, video or title slide"

    kinds = {str(safe_shape_type(s)).split()[0] for s in slide.shapes}
    if "TABLE" in kinds or "GRAPHIC" in kinds:
        return "contains a table"

    if any(safe_shape_type(s) != TEXT_BOX_TYPE for s in content):
        return "text sits inside shapes — looks like a diagram or callout"

    multi = [s for s in content if len(_para_texts(s)) >= 2]
    if len(multi) > 1:
        return "several multi-paragraph boxes — no single clear point list"
    if len(content) <= 2 and not multi:
        return "only a heading — no points to lay out"
    if len(content) >= 3 and not multi:
        # points would have to come from separate boxes; say why that failed
        fonts = Counter(round(_shape_font(s)) for s in content)
        biggest = max(fonts.values())
        if biggest < 2:
            return "content is one block of prose, not a list of points"
        item_font = min(f for f, c in fonts.items() if c == biggest)
        items = [s for s in content if round(_shape_font(s)) == item_font]
        tops = sorted(s.top or 0 for s in items)
        if any(tops[i + 1] - tops[i] < SAME_ROW_EMU for i in range(len(tops) - 1)):
            return "points sit side by side, not stacked as a list"
        lefts = [s.left or 0 for s in items]
        if max(lefts) - min(lefts) > ALIGN_TOL_EMU:
            return "points are scattered across the slide, not one column"
        return "surrounding text could not be read as a title/subtitle"
    if len(content) == 2 and multi:
        return "a single paragraph of prose, not a point list"
    return "no repeating point boxes found"


def _extract_single_body(body, others):
    """Heading / subheading / bullets for a Type E slide."""
    ordered = sorted(others, key=lambda s: (s.top or 0, s.left or 0))
    heading = ordered[0].text_frame.text.strip()
    subheading = None
    if len(ordered) >= 2 and (ordered[1].top or 0) <= (body.top or 0):
        subheading = ordered[1].text_frame.text.strip()
    return heading, subheading, _para_texts(body)


def _shape_sig(s):
    """Identity of a shape by kind + position/size on the slide, quantised
    to 1000-EMU (~0.03 mm) buckets: exporters can drift chrome by an EMU or
    two between slides, silently demoting repeated chrome from decoration
    to content (seen live in the ppsu1 fork; ported here)."""
    def q(v):
        return round((v or 0) / 1000)
    return (str(safe_shape_type(s)), q(s.top), q(s.left),
            q(s.width), q(s.height))


def _decoration_signatures(prs):
    """Shapes that repeat at the SAME spot on most slides of the deck.

    These are the deck's own chrome - the header rule, the logo, the course-code
    footer - not slide content, so they can be ignored when deciding whether a

    slide is a plain list of points.
    """
    n = len(prs.slides)
    counts = Counter()
    for slide in prs.slides:
        for s in slide.shapes:
            counts[_shape_sig(s)] += 1
    threshold = max(2, n * 0.5)
    return {k for k, c in counts.items() if c >= threshold}


AUTO_SHAPE_TYPE = 1  # MSO_SHAPE_TYPE.AUTO_SHAPE


def _is_plain_text_shape(s):
    """A plain text container: a text box, or an auto-shape with NO fill and NO
    border (so it reads as just text on the slide). A filled/outlined shape is
    a styled card or callout, not a bare point.

    A type-None shape (one python-pptx cannot classify — see
    safe_shape_type) that CARRIES text is judged exactly like the
    auto-shape case: fill-less/border-less means it reads as plain text."""
    st = safe_shape_type(s)
    if st == TEXT_BOX_TYPE:
        return True
    if st is None and not getattr(s, "has_text_frame", False):
        return False
    if st not in (AUTO_SHAPE_TYPE, None):
        return False          # groups, pictures, tables, etc.
    try:                      # MSO_FILL_TYPE.BACKGROUND (5) == transparent
        ft = s.fill.type
        if ft is not None and int(ft) != 5:
            return False
    except Exception:
        pass
    try:
        if s.line.width and int(s.line.width) > 0:
            return False
    except Exception:
        pass
    return True


def _items_by_font(content, slide_shapes, decor):
    """Type D raw slide: split shapes into repeating points vs surrounding text
    using font size alone.

    Used when the points share no number badge and no common shape name, so the
    only thing tying them together is that they are all set at the same size
    while the title / subtitle / footer are set at different ones.

    Returns (item_shapes, other_shapes), or None when there is no convincing
    group of same-sized point boxes (a title slide, a photo slide, a diagram),
    in which case the caller leaves the slide alone. Three or more same-sized
    boxes qualify on their own; a pair qualifies only with a title AND subtitle
    above it.
    """
    if len(content) < 4:  # need heading + subtitle + at least 2 points
        return None

    sized = [(s, round(_shape_font(s))) for s in content]
    counts = Counter(f for _, f in sized)
    biggest = max(counts.values())
    if biggest < 2:
        return None
    # on a tie prefer the smaller size: point text is smaller than a heading
    item_font = min(f for f, c in counts.items() if c == biggest)

    items = [s for s, f in sized if f == item_font]
    others = [s for s, f in sized if f != item_font]
    if not others:  # nothing left to serve as the heading
        return None

    # A group of only TWO same-sized boxes is weak evidence by itself (it could
    # be a section divider), so trust it only when there is an unambiguous title
    # AND subtitle above - two distinct larger font sizes. THREE or more
    # same-sized points remain sufficient on their own, exactly as before, so
    # decks that already worked are unaffected.
    if biggest < 3:
        bigger = {f for _, f in sized if f > item_font}
        if len(bigger) < 2:
            return None

    # ---- guards: only a plain vertical stack of point boxes qualifies. ----
    # Anything else (a diagram, a two-column layout, a slide with an extra
    # caption) is left completely untouched, so the original slide is carried
    # into the final deck exactly as it was rather than being rebuilt and
    # losing its shapes, arrows or stray text.

    # 1) Every point must be a PLAIN text container — a text box, or a fill-less
    #    border-less auto-shape (visually just text). A filled/outlined shape is
    #    a styled card or callout, so a slide built from those is left as-is.
    #    And a mere PAIR of auto-shapes is too weak (likely a callout+body
    #    layout), so auto-shape points are only trusted when there are 3+.
    if any(not _is_plain_text_shape(s) for s in items):
        return None
    if len(items) < 3 and any(safe_shape_type(s) != TEXT_BOX_TYPE for s in items):
        return None

    # 2) points must be stacked vertically. Two boxes sharing a row are a
    #    side-by-side layout (e.g. "Advantages | Disadvantages"), not a list.
    tops = sorted(s.top or 0 for s in items)
    if any(tops[i + 1] - tops[i] < SAME_ROW_EMU for i in range(len(tops) - 1)):
        return None

    # 3) every other text shape must be explainable as the title/subtitle above
    #    the points, or a small footer below them. A leftover block (a diagram
    #    caption, an extra note) would be silently dropped, so bail out.
    item_top, item_bottom = tops[0], tops[-1]
    for s in others:
        f = round(_shape_font(s))
        top = s.top or 0
        if top <= item_top and f > item_font:
            continue  # title / subtitle above the points
        if top > item_bottom and f < item_font:
            continue  # small footer underneath
        return None   # unexplained content -> keep the slide as it is

    # 3b) a real list is one column: its points share a left edge and a width.
    #     Diagram captions ("1. Single", "2. Multilevel", "3. Hierarchical") sit
    #     at different x positions with different widths, so they fail here.
    lefts = [s.left or 0 for s in items]
    widths = [s.width or 0 for s in items]
    if max(lefts) - min(lefts) > ALIGN_TOL_EMU:
        return None
    if max(widths) - min(widths) > WIDTH_TOL_EMU:
        return None

    # 4) every picture / group / auto-shape must be explainable too. Per-point
    #    decoration (a bullet dot beside a point, a card behind it) lines up
    #    vertically with one of the points and is safe to replace. A graphic
    #    that lines up with nothing is a real diagram - an arrow, a flow chart,
    #    an illustration - and rebuilding the slide would destroy it, so the
    #    slide is left exactly as it is instead.
    spans = [((s.top or 0), (s.top or 0) + (s.height or 0)) for s in items]
    for s in slide_shapes:
        if s.has_text_frame and s.text_frame.text.strip():
            continue                      # text already checked above
        if _shape_sig(s) in decor:
            continue                      # deck chrome (logo, rule, footer)
        top = s.top or 0
        bottom = top + (s.height or 0)
        if any(top < hi and bottom > lo for lo, hi in spans):
            continue                      # sits alongside one of the points
        return None                       # unexplained graphic -> keep as is

    return items, others


def extract_raw_slides(raw_path: Path, reasons=None):
    """Pull {heading, subheading, bullets[]} out of each raw slide.

    Pass a dict as `reasons` to receive {slide_number: why_it_was_skipped} for
    every slide that could not be turned into a design — used to explain each
    kept slide in the UI.

    Handles two raw-deck styles automatically:

    * Type A — all bullets live in a single text box (one paragraph each).
      Detected by explicit Heading / Subheading / BulletBox shape names, or by
      a single multi-paragraph text box in the positional fallback.
    * Type B — every bullet is its own separate text box (4 points -> 4 boxes),
      usually with a numbered badge shape beside each. Detected when no single
      multi-line bullet box is present; see `_extract_multibox`.

    Both styles yield the same {heading, subheading, bullets, count} record, so
    everything downstream (layout matching, adaptation, animation copy) is
    unchanged.
    """
    prs = Presentation(str(raw_path))
    # shapes repeating in the same spot across the deck are its chrome, not
    # slide content; worked out once so Type D can ignore them below
    decor = _decoration_signatures(prs)
    slides = []
    for idx, slide in enumerate(prs.slides, 1):
        named = {shp.name: shp for shp in slide.shapes if shp.has_text_frame}
        heading = subheading = bullets = None
        pairs = None
        skip_content = [s for s in slide.shapes
                        if s.has_text_frame and s.text_frame.text.strip()
                        and not RAW_BADGE_RE.match(s.text_frame.text.strip())]

        # Type A: an explicit "Heading" plus ONE box holding every point as its
        # own paragraph. Decks spell that box either "BulletBox" or "Content",
        # and the subheading either "Subheading" or "SubHead".
        body_name = next((n for n in ("BulletBox", "Content") if n in named), None)
        if "Heading" in named and body_name is not None:
            bullets = _para_texts(named[body_name])
            # "Content" decks also use that box for a single prose paragraph on
            # some slides; one paragraph is not a list of points, so leave those
            # slides untouched. ("BulletBox" decks behave exactly as before.)
            if body_name == "Content" and len(bullets) < 2:
                bullets = None
            else:
                heading = named["Heading"].text_frame.text.strip()
                sub = named.get("Subheading") or named.get("SubHead")
                subheading = sub.text_frame.text.strip() if sub else None
        else:
            text_shapes = [s for s in slide.shapes
                           if s.has_text_frame and s.text_frame.text.strip()]
            # Number badges ("1", "2", "01" ...) mark the repeating items in a
            # Type B deck. A slide is only treated as a redesignable content
            # slide when that numbered pattern is present (>= 3 badges); every
            # other slide -- title / section slides, image / table / diagram
            # slides, and badge-less layouts -- is not recognised here and is
            # left for build() to carry into the final deck unchanged.
            badges = [s for s in text_shapes
                      if RAW_BADGE_RE.match(s.text_frame.text.strip())]
            content = [s for s in text_shapes if s not in badges]
            if len(badges) >= 3 and len(content) >= 2:
                content.sort(key=lambda s: (s.top or 0, s.left or 0))
                heading, subheading, bullets = _extract_multibox(content)
            else:
                # Type C: a cluster of >= 3 text boxes sharing the exact same
                # shape name (e.g. "PointText") are the repeating items, each an
                # optional label + description. This name-repeat signal is unique
                # to these decks, so Type A / B slides are unaffected.
                name_counts = Counter(s.name for s in text_shapes)
                repeated = {n for n, c in name_counts.items() if c >= 3}
                # a deck that explicitly names its point boxes is telling us
                # what they are, so accept a pair of them (see the constant)
                if name_counts.get(EXPLICIT_POINT_NAME, 0) >= 2:
                    repeated.add(EXPLICIT_POINT_NAME)
                item_shapes = other_shapes = None
                if repeated:
                    item_name = max(repeated, key=lambda n: name_counts[n])
                    candidates = [s for s in text_shapes if s.name == item_name]
                    min_items = 2 if item_name == EXPLICIT_POINT_NAME else 3
                    # a repeated name is not proof of a list: a diagram's boxes
                    # are often all named alike too (e.g. 8 "DiagBox" shapes
                    # holding ClassA / ClassB / ClassC side by side). Only accept
                    # it when those shapes really form a column of point boxes.
                    if _is_point_column(candidates, min_items):
                        item_shapes = candidates
                        other_shapes = [s for s in text_shapes
                                        if s.name != item_name]
                if item_shapes is None:
                    # Type D: the points carry neither a number badge nor a shared
                    # shape name -- every point is its own uniquely-named box (e.g.
                    # "TextBox 8/10/12/14"). What they DO share is one font size,
                    # with a bigger title/subtitle above them. Reached only after
                    # Types A/B/C have all failed, so existing decks are untouched.
                    grouped = _items_by_font(content, slide.shapes, decor)
                    if grouped is not None:
                        item_shapes, other_shapes = grouped

                if item_shapes is not None:
                    heading, subheading, bullets, pairs = _extract_named_items(
                        item_shapes, other_shapes)
                else:
                    # Type E: no repeating boxes at all — the points live as
                    # separate paragraphs inside ONE text box, under generic
                    # names. Last resort, so every other deck is untouched.
                    single = _single_body_box(content, slide.shapes, decor)
                    if single is None:
                        if reasons is not None:
                            reasons[idx] = _skip_reason(slide, skip_content)
                        continue
                    heading, subheading, bullets = _extract_single_body(*single)

        if not bullets:
            if reasons is not None:
                reasons[idx] = _skip_reason(slide, skip_content)
            continue

        slides.append({
            "index": idx,
            "heading": heading,
            "subheading": subheading,
            "bullets": bullets,
            "count": len(bullets),
            "pairs": pairs,
        })
    return slides


def short_label(sentence: str, max_words: int = 4) -> str:
    """Derive a short label for two-field template items (label + description)."""
    words = sentence.split()
    label = " ".join(words[:max_words])
    return label.rstrip(",.;:")


SPLIT_DELIMS = [", and ", ", ", " and ", " — ", "; "]


def _find_split_point(sentence: str):
    """Find the delimiter closest to the sentence's midpoint to split on."""
    candidates = []
    for delim in SPLIT_DELIMS:
        start = 0
        while True:
            idx = sentence.find(delim, start)
            if idx == -1:
                break
            candidates.append(idx + len(delim))
            start = idx + 1
    if not candidates:
        return None
    mid = len(sentence) / 2
    candidates.sort(key=lambda p: abs(p - mid))
    return candidates[0]


def adapt_bullets(bullets, target_count):
    """Reshape a bullet list to exactly target_count items by merging (at most
    2 bullets per item) or splitting (at most 1 split per bullet) on natural
    delimiters. Returns None if the gap is too large to bridge cleanly."""
    n = len(bullets)
    if n == target_count:
        return list(bullets), False

    if target_count < n:
        # merging is always possible (just join text); heavier merges cost
        # more (see caller's cost ranking) but are never structurally blocked
        base, rem = divmod(n, target_count)
        sizes = [base + 1] * rem + [base] * (target_count - rem)
        result, idx = [], 0
        for size in sizes:
            result.append("; ".join(bullets[idx:idx + size]))
            idx += size
        return result, True

    extra = target_count - n  # number of splits needed
    splittable = [(i, _find_split_point(b)) for i, b in enumerate(bullets)]
    splittable = [(i, p) for i, p in splittable if p is not None]
    if len(splittable) < extra:
        return None
    splittable.sort(key=lambda x: -len(bullets[x[0]]))
    chosen = dict(splittable[:extra])
    result = []
    for i, b in enumerate(bullets):
        if i in chosen:
            p = chosen[i]
            first = b[:p].rstrip(" ,;—").strip()
            second = b[p:].strip()
            result.append(first)
            result.append(second)
        else:
            result.append(b)
    return result, True


def find_by_id(shapes, sid):
    for shp in shapes:
        if shp.shape_id == sid:
            return shp
        if safe_shape_type(shp) == 6:  # GROUP
            found = find_by_id(shp.shapes, sid)
            if found is not None:
                return found
    return None


# Breathing space inside every box we write text into (EMU).
# 91440 = 0.10 inch left/right; 45720 = 0.05 inch top/bottom — PowerPoint's own
# defaults, applied uniformly so no box has text glued to its edge.
TEXT_INSET_LR = Emu(91440)
TEXT_INSET_TB = Emu(45720)


def set_shape_text(shape, text, *, font_size=None, font_name=None):
    """Replace a shape's text while applying the requested font settings, and
    normalise the text's alignment INSIDE the box (the box itself is never
    moved or resized):

    * extra paragraphs are REMOVED (not just emptied) — a leftover empty
      paragraph is an invisible blank line that pushes the text off-centre;
    * the text is vertically centred (equal space above and below);
    * uniform internal margins give the text breathing space;
    * PowerPoint's "shrink text on overflow" autofit is disabled, so the font
      size set here is exactly what renders — text is never silently shrunk.
    """
    tf = shape.text_frame
    para = tf.paragraphs[0]
    if not para.runs:
        para.text = text
    else:
        para.runs[0].text = text
        for extra in para.runs[1:]:
            extra.text = ""
    # DELETE any additional paragraphs beyond the first — an empty <a:p> still
    # renders as a blank line and would break vertical centring
    for p in tf.paragraphs[1:]:
        p._p.getparent().remove(p._p)

    if para.runs:
        run = para.runs[0]
        if font_size is not None:
            run.font.size = Pt(font_size)
        if font_name is not None:
            run.font.name = font_name

    # ---- text alignment inside the (untouched) box ----
    try:
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE          # never shrink the text
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE     # equal top/bottom space
        tf.margin_left = TEXT_INSET_LR
        tf.margin_right = TEXT_INSET_LR
        tf.margin_top = TEXT_INSET_TB
        tf.margin_bottom = TEXT_INSET_TB
    except Exception:
        pass  # alignment is best-effort; the text itself is already in place


def _line_height(font_pt):
    return int(font_pt * 1.22 * 12700)           # line height incl. leading, EMU


def _needed_text_height(text, font_pt, box_w):
    """Estimated height (EMU) `text` needs at `font_pt` inside a box of width
    `box_w`, using average glyph metrics and a greedy word-wrap simulation.
    Returns 0 when it cannot be estimated."""
    avail_w = int(box_w or 0) - int(TEXT_INSET_LR) * 2
    if avail_w <= 0 or not text:
        return 0
    char_w = font_pt * 0.52 * 12700              # avg glyph width, EMU
    chars_per_line = max(1, int(avail_w / char_w))
    lines, cur = 1, 0
    for word in text.split():
        need = len(word) + (1 if cur else 0)
        if cur + need <= chars_per_line:
            cur += need
        else:
            lines += 1
            cur = len(word)
    return lines * _line_height(font_pt)


def text_overflows(text, font_pt, box_w, box_h, tolerance_lines=1.0):
    """True when `text` needs more than `tolerance_lines` extra lines beyond
    the box. Template boxes are routinely authored snug (a one-line subtitle in
    a box exactly one line tall renders fine), so a little nominal overflow is
    normal; callers pick how strict to be."""
    try:
        avail_h = int(box_h or 0) - int(TEXT_INSET_TB) * 2
        needed = _needed_text_height(text, font_pt, box_w)
        if avail_h <= 0 or needed <= 0:
            return False
        return needed > avail_h + tolerance_lines * _line_height(font_pt)
    except Exception:
        return False


def estimate_overflow(shape, text, font_pt, tolerance_lines=1.0):
    """Shape-based wrapper around text_overflows."""
    try:
        return text_overflows(text, font_pt, shape.width, shape.height,
                              tolerance_lines)
    except Exception:
        return False


# ---------------------------------------------------------------------------
# Pure python-pptx slide duplication / deletion — no zip unpacking, no
# external scripts. This is what previously depended on sandbox-only helper
# scripts (add_slide.py / clean.py); it's now self-contained.
# ---------------------------------------------------------------------------

def duplicate_slide(prs, index):
    """Clone slide `index` (0-based) of `prs`, appending the copy at the end
    of the deck. Shapes, background, and the relationships needed to render
    them (images, embedded media, external hyperlinks) are copied, and shape
    IDs are preserved so callers can keep looking shapes up by the IDs
    recorded in the catalog. Speaker notes, comments, and any internal
    hyperlink pointing at another slide are deliberately NOT carried over —
    they belong to the original deck, not the rebuilt one. Returns the new
    Slide.
    """
    source = prs.slides[index]
    dest = prs.slides.add_slide(source.slide_layout)

    # add_slide() clones placeholders from the layout; drop them since we're
    # about to copy the source slide's actual shapes wholesale
    for shp in list(dest.shapes):
        shp._element.getparent().remove(shp._element)

    # copy only the relationships a shape actually needs to render (images,
    # embedded media/objects, external hyperlinks) — skip the layout link
    # (dest already has its own), speaker notes, comments, and any internal
    # link to another slide (that slide won't necessarily still exist once
    # the deck is rebuilt, and dragging it along would pull in unrelated
    # original content). Remember old rId -> new rId so copied shape XML can
    # be repointed.
    rid_map = {}
    for rId, rel in list(source.part.rels.items()):
        if rel.reltype in (RT.SLIDE_LAYOUT, RT.NOTES_SLIDE):
            continue
        if not rel.is_external and isinstance(rel.target_part, SlidePart):
            continue  # internal "jump to slide" link — don't carry it over
        if rel.is_external:
            new_rid = dest.part.relate_to(rel.target_ref, rel.reltype, is_external=True)
        else:
            new_rid = dest.part.relate_to(rel.target_part, rel.reltype)
        rid_map[rId] = new_rid

    # deep-copy every shape, rewriting any r:id/r:embed/r:link references;
    # if a reference wasn't carried over above (e.g. a skipped internal
    # slide link), drop the attribute entirely rather than leave it dangling
    for shp in source.shapes:
        new_el = copy.deepcopy(shp._element)
        for el in new_el.iter():
            for attr, val in list(el.attrib.items()):
                if attr.startswith(f"{{{R_NS}}}"):
                    if val in rid_map:
                        el.set(attr, rid_map[val])
                    else:
                        del el.attrib[attr]
        dest.shapes._spTree.append(new_el)

    # best-effort background copy (slide-specific background overrides)
    try:
        src_cSld = source._element.find(qn("p:cSld"))
        src_bg = src_cSld.find(qn("p:bg")) if src_cSld is not None else None
        if src_bg is not None:
            dest_cSld = dest._element.find(qn("p:cSld"))
            existing_bg = dest_cSld.find(qn("p:bg"))
            if existing_bg is not None:
                dest_cSld.remove(existing_bg)
            dest_cSld.insert(0, copy.deepcopy(src_bg))
    except Exception:
        pass

    # copy the animation timeline (p:timing) so entrance / emphasis / exit
    # animations authored on the template slide are preserved on the clone.
    # The timing tree targets shapes by their shape id (spid), and the shape
    # copy above keeps those ids intact, so the references stay valid. Any
    # relationship references inside it (e.g. an audio/video trigger) are
    # repointed through the same rid_map used for the shapes; a reference we
    # didn't carry over is dropped rather than left dangling. p:timing must be
    # the last child of p:sld, so it's appended at the end of the slide.
    src_timing = source._element.find(qn("p:timing"))
    if src_timing is not None:
        new_timing = copy.deepcopy(src_timing)
        for el in new_timing.iter():
            for attr, val in list(el.attrib.items()):
                if attr.startswith(f"{{{R_NS}}}"):
                    if val in rid_map:
                        el.set(attr, rid_map[val])
                    else:
                        del el.attrib[attr]
        existing_timing = dest._element.find(qn("p:timing"))
        if existing_timing is not None:
            dest._element.remove(existing_timing)
        dest._element.append(new_timing)

    return dest


def delete_slide(prs, index):
    """Remove slide `index` (0-based) from the deck entirely (both from the
    slide order and the underlying relationship, so the part is dropped on
    save)."""
    sldIdLst = prs.slides._sldIdLst
    id_elem = list(sldIdLst)[index]
    rId = id_elem.get(qn("r:id"))
    sldIdLst.remove(id_elem)
    prs.part.drop_rel(rId)


def _blank_layout(prs):
    """Pick the layout with the fewest placeholders (closest to blank) to host
    an imported as-is slide, so the host layout contributes no stray shapes."""
    layouts = [lo for master in prs.slide_masters for lo in master.slide_layouts]
    layouts.sort(key=lambda lo: len(lo.placeholders))
    return layouts[0] if layouts else prs.slides[0].slide_layout


def _clone_part(dest_prs, part):
    """Copy a foreign part's bytes into dest_prs under a fresh, non-colliding
    partname. dest may already hold e.g. /ppt/media/image1.png with different
    content, so we never reuse the original partname."""
    import posixpath
    from pptx.opc.package import Part

    pkg = dest_prs.part.package
    base = str(part.partname)
    dirname = posixpath.dirname(base)
    stem = re.sub(r"\d+$", "", posixpath.splitext(posixpath.basename(base))[0]) or "part"
    ext = posixpath.splitext(base)[1]
    new_partname = pkg.next_partname("%s/%s%%d%s" % (dirname, stem, ext))
    return Part(new_partname, part.content_type, pkg, part.blob)


def copy_slide_asis(dest_prs, source_slide, dest_layout, drop_logo_sigs=frozenset()):
    """Append `source_slide` (from another presentation) into dest_prs exactly
    as it is — shapes, background, its own animations, and any image/media
    parts. Used to carry raw slides that aren't redesigned (image, title,
    table, diagram slides, or content that matched no template) into the final
    deck untouched. Image/media parts are re-packaged under fresh partnames to
    avoid clashing with the template deck's own media.

    `drop_logo_sigs` holds the shape signatures of the RAW deck's own repeating
    logo. The design template's master already draws a logo, in the position the
    rebuilt slides use, so copying the raw one as well would show two — and the
    raw one sits at a slightly different height. Skipping it keeps the logo in
    the same place on every slide of the finished deck.
    """
    dest = dest_prs.slides.add_slide(dest_layout)
    # add_slide() clones the layout's placeholders; drop them, we copy the
    # source slide's real shapes wholesale
    for shp in list(dest.shapes):
        shp._element.getparent().remove(shp._element)

    # copy the relationships the shapes render from (images, media, external
    # hyperlinks); skip the layout/notes links and any internal slide link
    rid_map = {}
    for rId, rel in list(source_slide.part.rels.items()):
        if rel.reltype in (RT.SLIDE_LAYOUT, RT.NOTES_SLIDE):
            continue
        if not rel.is_external and isinstance(rel.target_part, SlidePart):
            continue
        if rel.is_external:
            new_rid = dest.part.relate_to(rel.target_ref, rel.reltype, is_external=True)
        else:
            cloned = _clone_part(dest_prs, rel.target_part)
            new_rid = dest.part.relate_to(cloned, rel.reltype)
        rid_map[rId] = new_rid

    for shp in source_slide.shapes:
        # skip the raw deck's own logo — the design master already draws one,
        # and at the height the rebuilt slides use, so keeping both would show
        # two logos at slightly different heights
        if _shape_sig(shp) in drop_logo_sigs:
            continue
        new_el = copy.deepcopy(shp._element)
        for el in new_el.iter():
            for attr, val in list(el.attrib.items()):
                if attr.startswith(f"{{{R_NS}}}"):
                    if val in rid_map:
                        el.set(attr, rid_map[val])
                    else:
                        del el.attrib[attr]
        dest.shapes._spTree.append(new_el)

    # background override, if the source slide has one
    try:
        src_cSld = source_slide._element.find(qn("p:cSld"))
        src_bg = src_cSld.find(qn("p:bg")) if src_cSld is not None else None
        if src_bg is not None:
            dest_cSld = dest._element.find(qn("p:cSld"))
            existing_bg = dest_cSld.find(qn("p:bg"))
            if existing_bg is not None:
                dest_cSld.remove(existing_bg)
            dest_cSld.insert(0, copy.deepcopy(src_bg))
    except Exception:
        pass

    # keep any animations the raw slide already had
    src_timing = source_slide._element.find(qn("p:timing"))
    if src_timing is not None:
        new_timing = copy.deepcopy(src_timing)
        for el in new_timing.iter():
            for attr, val in list(el.attrib.items()):
                if attr.startswith(f"{{{R_NS}}}"):
                    if val in rid_map:
                        el.set(attr, rid_map[val])
                    else:
                        del el.attrib[attr]
        dest._element.append(new_timing)

    return dest


def _style_title_slide(slide):
    """Apply the fixed title-slide typography requested for the 2nd slide.

    The text boxes are taken top-to-bottom: the orange badge (28pt), the
    subject name (36pt), then the unit box whose first line is the unit number
    (32pt) and whose remaining line is the unit name (28pt). All set to Aptos.
    Only font size/name are changed; text and everything else are left as-is.
    """
    boxes = [s for s in slide.shapes
             if s.has_text_frame and s.text_frame.text.strip()]
    boxes.sort(key=lambda s: (s.top or 0, s.left or 0))

    def restyle(paragraph, size):
        for r in paragraph.runs:
            r.font.size = Pt(size)
            r.font.name = "Aptos"

    if len(boxes) >= 1:  # orange badge
        for p in boxes[0].text_frame.paragraphs:
            restyle(p, 28)
    if len(boxes) >= 2:  # subject name
        for p in boxes[1].text_frame.paragraphs:
            restyle(p, 36)
    if len(boxes) >= 3:  # unit box: number then name
        unit_paras = [p for p in boxes[2].text_frame.paragraphs if p.text.strip()]
        for i, p in enumerate(unit_paras):
            restyle(p, 32 if i == 0 else 28)


def _design_of(entry):
    """Design-family id for an entry: its style fingerprint if the catalog has
    one, else the slide file (so older catalogs still work, just ungrouped)."""
    return entry.get("design_key") or entry["slide_file"]


def _choose_pool(rs, reserved_templates, general_templates):
    """Eligible designs for this raw slide, plus whether it is a pinned section.

    A reserved-section slide (heading == "Summary" / "Learning Objectives" ...)
    is pinned to the dedicated designs carrying the same title; if the template
    has none, it falls back to the general pool. Every other slide uses the
    general pool, which never contains the reserved designs.

    Returns (pool, is_reserved)."""
    raw_title = _norm_title(rs["heading"])
    if raw_title in RESERVED_TITLES:
        pinned = [e for e in reserved_templates if e["_title_norm"] == raw_title]
        if pinned:
            return pinned, True
    return general_templates, False


def _select_design(rs, pool, last_used, exact_first=False):
    """Pick the best design in `pool` for one raw slide, breaking final ties
    RANDOMLY (so there's never a fixed catalog-order pattern).

    exact_first=True (used for all content): match the point count EXACTLY
    first — a 3-point slide takes a 3-item design, a 4-point one a 4-item
    design — so a slide is never inflated into a bigger design (which would
    fabricate a point by splitting) while a same-count design exists. Among
    designs of that count, the least-recently-used wins, so nothing repeats
    until every design of that count has been used; the rotation state is
    shared across sibling decks (unit1.1 / 1.2 / 1.3) so the spread continues
    between them. Adapting to a different count happens only as a last resort,
    when the template has no design of the exact count.

    exact_first=False keeps the older least-recently-used-first behaviour and
    is retained only for callers that may want maximum design variety over an
    exact count match.

    Returns (entry, adapted_bullets, was_adapted) or None if nothing fits."""
    candidates = []
    for entry in pool:
        # Only a design with the EXACT number of items may be used. If the
        # template has no design for this point count (e.g. a 6-point slide
        # when the template stops at 5), the slide is left exactly as it is
        # rather than having its points merged or split to fit.
        if entry["item_count"] != rs["count"]:
            continue
        adapted = adapt_bullets(rs["bullets"], entry["item_count"])
        if adapted is None:
            continue
        adapted_bullets, was_adapted = adapted
        cost = abs(entry["item_count"] - rs["count"])
        lu = last_used.get(_design_of(entry), -1)
        rank = (cost, lu) if exact_first else (lu, cost)
        candidates.append((rank, entry, adapted_bullets, was_adapted))
    if not candidates:
        return None
    best_rank = min(c[0] for c in candidates)
    _, entry, adapted_bullets, was_adapted = random.choice(
        [c for c in candidates if c[0] == best_rank])
    return entry, adapted_bullets, was_adapted


KEEP_ORIGINAL = "__keep__"


def build(raw_path: Path, template_path: Path, catalog_path: Path, out_path: Path,
          usage_state=None, fixed=None):
    """Rebuild `raw_path` using `template_path`'s designs.

    `fixed` pins per-slide decisions so a deck can be reproduced exactly (needed
    to change ONE slide without disturbing the rest). It maps a raw slide number
    to either a design's slide_file (force that design) or KEEP_ORIGINAL (carry
    the slide over untouched). Slides not listed are chosen normally.
    """
    fixed = fixed or {}
    raw_catalog = json.loads(catalog_path.read_text())
    # accept both shapes: {"count": {...one entry...}} (legacy) and
    # {"count": [...multiple entries...]} (auto-detected, for rotation)
    all_templates = []
    for k, v in raw_catalog["patterns"].items():
        entries = v if isinstance(v, list) else [v]
        for e in entries:
            e.setdefault("item_count", len(e["items"]))
            all_templates.append(e)

    # Design-rotation state, keyed by design family so that look-alike layouts
    # (e.g. a 4-item and a 5-item version of the same chevron design) count as
    # ONE design. When the caller passes a shared dict, the rotation continues
    # across sibling decks (e.g. unit1.1 / 1.2 / 1.3), so a design used in one
    # unit isn't reused in the next until every other design has been cycled.
    if usage_state is None:
        usage_state = {}
    last_used = usage_state.setdefault("last_used", {})

    # Tag which template designs are "reserved" sections (Summary / Learning
    # Objectives ...) by their title text, so they can be pinned to matching
    # content and kept out of the general rotation.
    for e in all_templates:
        e["_title_norm"] = _norm_title(e.get("title_text"))
        e["_reserved"] = e["_title_norm"] in RESERVED_TITLES

    reserved_templates = [e for e in all_templates if e["_reserved"]]
    general_templates = [e for e in all_templates if not e["_reserved"]]

    skip_reasons = {}
    raw_slides = extract_raw_slides(raw_path, reasons=skip_reasons)

    by_file = {e["slide_file"]: e for e in all_templates}

    matched, unmatched = [], []
    for rs in raw_slides:
        pin = fixed.get(rs["index"]) or fixed.get(str(rs["index"]))
        if pin == KEEP_ORIGINAL:
            unmatched.append(rs)          # user chose to keep this slide as-is
            continue
        if pin and pin in by_file and by_file[pin]["item_count"] == rs["count"]:
            # reproduce / force a specific design (same point count only)
            entry = by_file[pin]
            adapted_bullets, was_adapted = adapt_bullets(rs["bullets"], entry["item_count"])
        else:
            pool, _ = _choose_pool(rs, reserved_templates, general_templates)
            # exact_first=True for ALL content: a 3-point slide must take a
            # 3-item design and never be inflated into a 4-item one (which
            # fabricates a point by splitting) as long as a same-count design
            # exists. Rotation then happens WITHIN the matching count group.
            choice = _select_design(rs, pool, last_used, exact_first=True)
            if choice is None:
                unmatched.append(rs)
                continue
            entry, adapted_bullets, was_adapted = choice
        usage_state["counter"] = usage_state.get("counter", 0) + 1
        last_used[_design_of(entry)] = usage_state["counter"]
        matched.append((rs, entry, adapted_bullets, was_adapted))

    prs = Presentation(str(template_path))
    original_count = len(prs.slides)
    raw_prs = Presentation(str(raw_path))
    total_raw = len(raw_prs.slides)
    blank_layout = _blank_layout(prs)

    # The raw deck stamps its own logo on every slide, and the design template's
    # master draws one too — at the height the rebuilt slides use. Carrying both
    # onto a kept slide shows two logos at slightly different heights, so drop
    # the raw one and let the design's (correctly placed) logo through. Only the
    # repeating picture is dropped: a picture used on a single slide is real
    # content and is always kept.
    drop_logo_sigs = frozenset()
    template_has_logo = any(
        "PICTURE" in str(safe_shape_type(s))
        for master in prs.slide_masters for s in master.shapes)
    if template_has_logo:
        raw_decor = _decoration_signatures(raw_prs)
        drop_logo_sigs = frozenset(sig for sig in raw_decor
                                   if "PICTURE" in sig[0])

    # analyze_template.py numbered slides via enumerate(prs.slides, 1) against
    # this same template file, so slide_file "slideN.xml" -> 0-based index N-1
    # here too (valid as long as the deck hasn't changed since the catalog
    # was built, which app.py guarantees by caching per file hash).
    def index_for(slide_file):
        return int(slide_file.replace("slide", "").replace(".xml", "")) - 1

    match_by_index = {rs["index"]: (rs, entry, adapted_bullets, was_adapted)
                      for rs, entry, adapted_bullets, was_adapted in matched}

    # 1) walk EVERY raw slide in its original order. A matched content slide is
    #    rebuilt from its template layout; every other slide (unmatched content,
    #    image / title / table / diagram slides) is carried over untouched, so
    #    nothing is ever dropped from the deck.
    built = []
    for raw_idx in range(1, total_raw + 1):
        if raw_idx in match_by_index:
            rs, entry, adapted_bullets, was_adapted = match_by_index[raw_idx]
            slide = duplicate_slide(prs, index_for(entry["slide_file"]))
            built.append(("matched", slide, rs, entry, adapted_bullets, was_adapted))
        else:
            slide = copy_slide_asis(prs, raw_prs.slides[raw_idx - 1], blank_layout,
                                    drop_logo_sigs=drop_logo_sigs)
            if raw_idx == 2:  # the 2nd slide: fixed title-slide typography
                _style_title_slide(slide)
            built.append(("kept", slide, None, None, None, None))

    # 2) drop every original template slide now that content has been cloned
    for i in range(original_count - 1, -1, -1):
        delete_slide(prs, i)

    # 3) fill in text on each rebuilt (matched) slide; kept slides are left alone.
    #    While filling, note every box whose text clearly needs more room than
    #    the box has at full font size — the text is left at full size (never
    #    shrunk), but the preview warns so the author can shorten it.
    overflow_by_slide = {}
    for kind, slide, rs, entry, adapted_bullets, was_adapted in built:
        if kind != "matched":
            continue
        too_long = []
        sub_shape = None
        title_text = rs["heading"]
        if rs["subheading"] and entry.get("subtitle_id") is None:
            title_text = f'{rs["heading"]} — {rs["subheading"]}'
        title_shape = find_by_id(slide.shapes, entry["title_id"])
        if title_shape is not None:
            set_shape_text(title_shape, title_text, font_size=40, font_name="Aptos")
            if estimate_overflow(title_shape, title_text, 40):
                too_long.append("title")

        if entry.get("subtitle_id") is not None:
            sub_shape = find_by_id(slide.shapes, entry["subtitle_id"])
            if sub_shape is not None:
                # Always write the subtitle box: the raw slide's subheading, or
                # BLANK it when the raw slide has none. Otherwise the reference
                # template's own subtitle text leaks into the finished slide.
                set_shape_text(sub_shape, rs["subheading"] or "",
                               font_size=28, font_name="Aptos")
                if rs["subheading"] and estimate_overflow(sub_shape, rs["subheading"], 28):
                    too_long.append("subtitle")

        # a raw slide may carry explicit (label, description) pairs per item
        # (Type C). Those align 1:1 with the bullets only when no merge/split
        # happened; otherwise fall back to deriving a label from the bullet.
        pairs = rs.get("pairs")
        use_pairs = bool(pairs) and not was_adapted and len(pairs) == len(adapted_bullets)

        items = entry["items"]
        for i, (bullet, item) in enumerate(zip(adapted_bullets, items)):
            pair = pairs[i] if use_pairs else None
            if isinstance(item, list):  # two-field: [label_id, desc_id]
                label_shape = find_by_id(slide.shapes, item[0])
                desc_shape = find_by_id(slide.shapes, item[1])
                label_text = pair[0] if pair else short_label(bullet)
                desc_text = pair[1] if pair else bullet
                if label_shape is not None:
                    set_shape_text(label_shape, label_text, font_size=22, font_name="Aptos")
                if desc_shape is not None:
                    set_shape_text(desc_shape, desc_text, font_size=22, font_name="Aptos")
                    if estimate_overflow(desc_shape, desc_text, 22):
                        too_long.append(f"point {i + 1}")
            else:
                shp = find_by_id(slide.shapes, item)
                if shp is not None:
                    set_shape_text(shp, bullet, font_size=22, font_name="Aptos")
                    if estimate_overflow(shp, bullet, 22):
                        too_long.append(f"point {i + 1}")
        if too_long:
            overflow_by_slide[rs["index"]] = too_long

    prs.save(str(out_path))

    report = {
        "matched": [{"raw_slide": rs["index"], "heading": rs["heading"],
                      "bullets": rs["count"], "template_slide": entry["slide_file"],
                      "template_items": entry["item_count"], "content_adapted": was_adapted}
                     for rs, entry, adapted_bullets, was_adapted in matched],
        "unmatched": [{"raw_slide": rs["index"], "heading": rs["heading"],
                        "bullets": rs["count"]} for rs in unmatched],
    }

    # One entry per slide of the finished deck, in order, explaining what
    # happened to it. The output deck is built by walking the raw slides in
    # order, so entry i lines up with preview image i.
    by_idx = {rs["index"]: (entry, rs) for rs, entry, _b, _a in matched}
    no_design = {rs["index"]: rs for rs in unmatched}
    explained = []
    for i in range(1, total_raw + 1):
        if i in by_idx:
            entry, rs = by_idx[i]
            row = {
                "raw_slide": i, "status": "redesigned",
                "heading": rs["heading"],
                "reason": f'{rs["count"]} points → {entry["item_count"]}-item design',
            }
            if i in overflow_by_slide:
                row["warning"] = ("text may be too long for its box: "
                                  + ", ".join(overflow_by_slide[i]))
            explained.append(row)
        elif i in no_design:
            rs = no_design[i]
            explained.append({
                "raw_slide": i, "status": "kept",
                "heading": rs["heading"],
                "reason": (f'{rs["count"]} points — the template has no '
                           f'{rs["count"]}-item design'),
            })
        else:
            explained.append({
                "raw_slide": i, "status": "kept", "heading": None,
                "reason": skip_reasons.get(i, "not recognised as a point list"),
            })
    report["slides"] = explained
    return report


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--raw", required=True, type=Path)
    ap.add_argument("--template", required=True, type=Path)
    ap.add_argument("--catalog", required=True, type=Path)
    ap.add_argument("--out", required=True, type=Path)
    args = ap.parse_args()
    report = build(args.raw, args.template, args.catalog, args.out)
    print(json.dumps(report, indent=2))


if __name__ == "__main__":
    main()
