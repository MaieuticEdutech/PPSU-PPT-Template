#!/usr/bin/env python3
"""
build_template.py — authors the PPSU design template deck.

The template is built ON TOP of a PPSU raw deck so it inherits the exact
PPSU slide-master chrome (navy banner, orange strip, decorative slashes,
PPSU logo on the white top-right panel). All original slides are removed
and replaced with design slides that analyze_template.py auto-detects:

* a "TitleText" box (white, bold, 24pt Calibri) sitting on the banner;
* N repeating "Item" shapes per design (the styled cards/pills that will
  receive the raw slide's points, 14pt Calibri);
* decorative accents and numbered chips that carry no real text (chip text
  like "01" is a badge and is ignored by the catalog).

Designs included: 2/3/4/5/6-point layouts (several looks per count for
rotation) plus dedicated reserved designs for "Learning Objectives" and
"Summary" slides (pinned by title, withheld from the general rotation).

Usage:  python build_template.py  (from the ppsu1 folder)
Writes: template/PPSU TEMPLATE.pptx
"""
import copy
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

HERE = Path(__file__).parent
RAW = HERE.parent / "PPSU RAW PPT" / "RAW PPT1.pptx"
OUT = HERE / "template" / "PPSU TEMPLATE.pptx"

# ---- PPSU brand palette (taken from the raw deck's master) -----------------
NAVY = RGBColor(0x0E, 0x28, 0x41)
DARK = RGBColor(0x0A, 0x1E, 0x30)
TEAL = RGBColor(0x15, 0x60, 0x82)
ORANGE = RGBColor(0xE9, 0x71, 0x32)
DEEP_ORANGE = RGBColor(0xF0, 0x4B, 0x22)
RED = RGBColor(0xB3, 0x11, 0x15)
CRIMSON = RGBColor(0xD8, 0x18, 0x1F)
GOLD = RGBColor(0xFB, 0xB2, 0x17)
BLUE = RGBColor(0x1D, 0x41, 0x81)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF4, 0xF7, 0xFA)       # near-white cool card fill
LIGHT_BLUE = RGBColor(0xEA, 0xF1, 0xF7)
LIGHT_WARM = RGBColor(0xFD, 0xF3, 0xE7)
LIGHT_NAVY = RGBColor(0xE8, 0xEE, 0xF5)
GRAY_LINE = RGBColor(0xD9, 0xDE, 0xE6)

ACCENTS = [ORANGE, TEAL, RED, GOLD, BLUE, DEEP_ORANGE]

FONT = "Calibri"

# ---- canvas geometry (EMU) — slide is 9144000 x 5143500 -------------------
X0, X1 = 380000, 8764000          # content left/right
CW = X1 - X0                      # 8384000 content width
Y0, Y1 = 880000, 4950000          # content top/bottom (below banner+strip)
CH = Y1 - Y0                      # 4070000 content height

TITLE_XY = (378920, 55000, 6500000, 570000)


def _no_shadow(shape):
    """Explicitly disable the theme's preset outer shadow on an autoshape."""
    spPr = shape._element.spPr
    el = spPr.find(qn("a:effectLst"))
    if el is None:
        el = spPr.makeelement(qn("a:effectLst"), {})
        spPr.append(el)
    else:
        for child in list(el):
            el.remove(child)


def style(shape, fill=None, line=None, line_w=0.75):
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w)
    _no_shadow(shape)


def set_text(shape, text, size, color, bold=False, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.MIDDLE, badge=False):
    tf = shape.text_frame
    # badges are tiny — no margins and no wrapping, or "01" breaks in two
    tf.word_wrap = not badge
    tf.vertical_anchor = anchor
    inset_lr = Emu(0) if badge else Emu(91440)
    inset_tb = Emu(0) if badge else Emu(45720)
    tf.margin_left = inset_lr
    tf.margin_right = inset_lr
    tf.margin_top = inset_tb
    tf.margin_bottom = inset_tb
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run() if not p.runs else p.runs[0]
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = FONT
    run.font.color.rgb = color


def add_shape(slide, prst, x, y, w, h, name, fill=None, line=None, line_w=0.75):
    s = slide.shapes.add_shape(prst, Emu(x), Emu(y), Emu(w), Emu(h))
    s.name = name
    style(s, fill, line, line_w)
    return s


def add_title(slide, text="Slide Title"):
    x, y, w, h = TITLE_XY
    box = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    box.name = "TitleText"
    set_text(box, text, 24, WHITE, bold=True)
    return box


def add_chip(slide, i, x, y, d, color, shape=MSO_SHAPE.OVAL, text_color=WHITE):
    """Small numbered badge ("01", "02"...) — decorative; the catalog ignores
    pure-number text, so chips never count as content."""
    c = add_shape(slide, shape, x, y, d, d, f"Chip {i + 1}", fill=color)
    set_text(c, "%02d" % (i + 1), 12, text_color, bold=True,
             align=PP_ALIGN.CENTER, badge=True)
    return c


ITEM_SAMPLE = "Your point text appears here in this layout"


def add_item(slide, i, x, y, w, h, fill, text_color, prst=MSO_SHAPE.ROUNDED_RECTANGLE,
             line=None, align=PP_ALIGN.LEFT, radius=None):
    s = add_shape(slide, prst, x, y, w, h, f"Item {i + 1}", fill=fill, line=line)
    if radius is not None:
        try:
            s.adjustments[0] = radius
        except Exception:
            pass
    set_text(s, ITEM_SAMPLE, 14, text_color, align=align)
    return s


def new_slide(prs, layout):
    slide = prs.slides.add_slide(layout)
    for shp in list(slide.shapes):     # drop cloned date/footer placeholders
        shp._element.getparent().remove(shp._element)
    return slide


def rows_geometry(n, gap):
    """(y, h) for n stacked rows filling the content area with `gap` between."""
    h = (CH - (n - 1) * gap) // n
    return [(Y0 + i * (h + gap), h) for i in range(n)]


def cols_geometry(n, gap, x0=X0, width=CW):
    w = (width - (n - 1) * gap) // n
    return [(x0 + i * (w + gap), w) for i in range(n)]


# ---------------------------------------------------------------------------
# the individual designs
# ---------------------------------------------------------------------------

def design_rows_chips(prs, layout, n, fill, text_color, chip_shape, gap=150000,
                      title="Slide Title", chip_colors=None, tab=False):
    """N full-width pill rows, each with a numbered chip on its left."""
    slide = new_slide(prs, layout)
    add_title(slide, title)
    chip_colors = chip_colors or ACCENTS
    for i, (y, h) in enumerate(rows_geometry(n, gap)):
        d = min(430000, h - 60000)
        cy = y + (h - d) // 2
        add_chip(slide, i, X0, cy, d, chip_colors[i % len(chip_colors)],
                 shape=chip_shape)
        ix = X0 + d + 130000
        add_item(slide, i, ix, y, X1 - ix, h, fill, text_color, radius=0.5)
        if tab:
            add_shape(slide, MSO_SHAPE.RECTANGLE, ix + 20000, y + h // 4,
                      45000, h // 2, f"Accent {i + 1}",
                      fill=chip_colors[i % len(chip_colors)])
    return slide


def design_grid(prs, layout, n, cols, fill, text_color, line=None,
                bar_colors=None, title="Slide Title"):
    """Cards in a grid (row-major), each with a colored top accent bar."""
    import math
    slide = new_slide(prs, layout)
    add_title(slide, title)
    rows = math.ceil(n / cols)
    gap = 200000
    row_h = (CH - (rows - 1) * gap) // rows
    bar_colors = bar_colors or ACCENTS
    k = 0
    for r in range(rows):
        remaining = n - k
        ncols = min(cols, remaining)
        xs = cols_geometry(ncols, gap)
        y = Y0 + r * (row_h + gap)
        for (x, w) in xs:
            if k >= n:
                break
            add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, 130000,
                      f"Accent {k + 1}", fill=bar_colors[k % len(bar_colors)])
            add_item(slide, k, x, y + 150000, w, row_h - 150000, fill,
                     text_color, line=line, radius=0.08)
            k += 1
    return slide


def design_columns(prs, layout, n, fill, text_color, title="Slide Title"):
    """N tall column cards with a numbered circle chip overlapping each top."""
    slide = new_slide(prs, layout)
    add_title(slide, title)
    gap = 220000
    d = 480000
    top = Y0 + d // 2
    for i, (x, w) in enumerate(cols_geometry(n, gap)):
        add_item(slide, i, x, top, w, Y1 - top, fill, text_color,
                 align=PP_ALIGN.CENTER, radius=0.07)
        add_chip(slide, i, x + (w - d) // 2, Y0, d,
                 ACCENTS[i % len(ACCENTS)], shape=MSO_SHAPE.OVAL)
    return slide


def design_arrows(prs, layout, n, title="Slide Title"):
    """N full-width arrow-pentagon bars in cycling brand colors, white text."""
    slide = new_slide(prs, layout)
    add_title(slide, title)
    colors = [NAVY, TEAL, ORANGE, RED, BLUE, DARK]
    for i, (y, h) in enumerate(rows_geometry(n, 140000)):
        add_item(slide, i, X0, y, CW - i * 260000, h, colors[i % len(colors)],
                 WHITE, prst=MSO_SHAPE.PENTAGON)
    return slide


def design_two_panels(prs, layout, title="Slide Title"):
    """2 items: two large side-by-side cards with colored top bars."""
    slide = new_slide(prs, layout)
    add_title(slide, title)
    y = Y0 + 300000
    h = Y1 - y
    for i, (x, w) in enumerate(cols_geometry(2, 300000)):
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, 150000,
                  f"Accent {i + 1}", fill=[NAVY, ORANGE][i])
        add_item(slide, i, x, y + 180000, w, h - 180000, LIGHT, NAVY,
                 line=GRAY_LINE, align=PP_ALIGN.CENTER, radius=0.05)
    return slide


def add_icon_tile(slide, i, x, y, d, color, prst=MSO_SHAPE.ROUNDED_RECTANGLE):
    """Empty colored tile named "Icon n" — the engine drops a matching icon
    PNG onto it at build time. Carries NO text, so the catalog's text
    clustering ignores it; analyze_template picks it up by its name."""
    t = add_shape(slide, prst, x, y, d, d, f"Icon {i + 1}", fill=color)
    try:
        t.adjustments[0] = 0.25
    except Exception:
        pass
    return t


def design_rows_icons(prs, layout, n, fill, text_color, gap=150000,
                      title="Slide Title", tile_colors=None):
    """N full-width pill rows, each with a colored icon tile on its left —
    the icon variant of design_rows_chips."""
    slide = new_slide(prs, layout)
    add_title(slide, title)
    tile_colors = tile_colors or ACCENTS
    for i, (y, h) in enumerate(rows_geometry(n, gap)):
        d = min(430000, h - 60000)
        cy = y + (h - d) // 2
        add_icon_tile(slide, i, X0, cy, d, tile_colors[i % len(tile_colors)])
        ix = X0 + d + 130000
        add_item(slide, i, ix, y, X1 - ix, h, fill, text_color, radius=0.5)
    return slide


def design_columns_icons(prs, layout, n, fill, text_color, title="Slide Title"):
    """N tall column cards with an icon tile overlapping each top — the icon
    variant of design_columns."""
    slide = new_slide(prs, layout)
    add_title(slide, title)
    gap = 220000
    d = 480000
    top = Y0 + d // 2
    for i, (x, w) in enumerate(cols_geometry(n, gap)):
        add_item(slide, i, x, top, w, Y1 - top, fill, text_color,
                 align=PP_ALIGN.CENTER, radius=0.07)
        add_icon_tile(slide, i, x + (w - d) // 2, Y0, d,
                      ACCENTS[i % len(ACCENTS)])
    return slide


def design_two_panels_icons(prs, layout, title="Slide Title"):
    """2 items: two large side-by-side cards, each with an icon tile centred
    on its top edge."""
    slide = new_slide(prs, layout)
    add_title(slide, title)
    d = 520000
    y = Y0 + 300000 + d // 2
    h = Y1 - y
    for i, (x, w) in enumerate(cols_geometry(2, 300000)):
        add_item(slide, i, x, y, w, h, LIGHT, NAVY,
                 line=GRAY_LINE, align=PP_ALIGN.CENTER, radius=0.05)
        add_icon_tile(slide, i, x + (w - d) // 2, y - d // 2, d,
                      [NAVY, ORANGE][i], prst=MSO_SHAPE.OVAL)
    return slide


def design_reserved_rows(prs, layout, n, title, card_fill, text_color,
                         chevron_fill, chevron_text):
    """Reserved-section design (Learning Objectives / Summary): rows with a
    chevron badge — visually distinct from the general rotation."""
    slide = new_slide(prs, layout)
    add_title(slide, title)
    gap = 160000
    for i, (y, h) in enumerate(rows_geometry(n, gap)):
        cw = 620000
        ch = min(430000, h - 60000)
        cy = y + (h - ch) // 2
        c = add_shape(slide, MSO_SHAPE.CHEVRON, X0, cy, cw, ch,
                      f"Chip {i + 1}", fill=chevron_fill)
        set_text(c, "%02d" % (i + 1), 12, chevron_text, bold=True,
                 align=PP_ALIGN.CENTER, badge=True)
        ix = X0 + cw + 130000
        add_item(slide, i, ix, y, X1 - ix, h, card_fill, text_color, radius=0.3)
    return slide


def main():
    prs = Presentation(str(RAW))

    # the "Blank" layout of the PPSU-branded master (used by every raw
    # content slide) — grab it from an existing slide before deleting them
    layout = prs.slides[1].slide_layout

    # remove every original slide
    sldIdLst = prs.slides._sldIdLst
    for id_elem in list(sldIdLst):
        rId = id_elem.get(qn("r:id"))
        sldIdLst.remove(id_elem)
        prs.part.drop_rel(rId)

    # ---- general designs ----
    design_two_panels(prs, layout)                                    # 2 items

    design_columns(prs, layout, 3, LIGHT_WARM, NAVY)                  # 3 items
    design_rows_chips(prs, layout, 3, LIGHT, NAVY, MSO_SHAPE.OVAL,
                      gap=230000)                                     # 3 items

    design_grid(prs, layout, 4, 2, WHITE, NAVY, line=GRAY_LINE)       # 4 items
    design_rows_chips(prs, layout, 4, LIGHT_BLUE, NAVY,
                      MSO_SHAPE.ROUNDED_RECTANGLE, gap=180000)        # 4 items
    design_columns(prs, layout, 4, LIGHT, NAVY)                       # 4 items

    design_rows_chips(prs, layout, 5, LIGHT, NAVY, MSO_SHAPE.OVAL,
                      gap=130000)                                     # 5 items
    design_grid(prs, layout, 5, 3, LIGHT_BLUE, NAVY,
                bar_colors=[TEAL, BLUE, NAVY, ORANGE, RED])           # 5 items
    design_arrows(prs, layout, 5)                                     # 5 items

    design_grid(prs, layout, 6, 3, WHITE, NAVY, line=GRAY_LINE)       # 6 items
    design_rows_chips(prs, layout, 6, LIGHT_WARM, NAVY,
                      MSO_SHAPE.ROUNDED_RECTANGLE, gap=110000,
                      chip_colors=[ORANGE, RED, GOLD, TEAL, BLUE, NAVY])  # 6

    # ---- icon designs (one per point count) ----
    # Each has an empty "Icon n" tile per item. The engine picks one of these
    # ONLY when every point of a raw slide gets a well-matched icon from the
    # icon service; otherwise it stays with the text-only designs above.
    design_two_panels_icons(prs, layout)                              # 2 items
    design_rows_icons(prs, layout, 3, LIGHT, NAVY, gap=230000)        # 3 items
    design_columns_icons(prs, layout, 4, LIGHT_WARM, NAVY)            # 4 items
    design_rows_icons(prs, layout, 5, LIGHT_BLUE, NAVY, gap=130000)   # 5 items
    design_rows_icons(prs, layout, 6, LIGHT, NAVY, gap=110000,
                      tile_colors=[ORANGE, RED, GOLD, TEAL, BLUE, NAVY])  # 6

    # ---- reserved section designs (pinned by title text) ----
    for n in (3, 4, 5):
        design_reserved_rows(prs, layout, n, "Learning Objectives",
                             LIGHT_NAVY, NAVY, GOLD, NAVY)
    for n in (3, 4, 5):
        design_reserved_rows(prs, layout, n, "Summary",
                             NAVY, WHITE, ORANGE, WHITE)

    OUT.parent.mkdir(exist_ok=True)
    prs.save(str(OUT))
    print("wrote", OUT, "with", len(prs.slides.__iter__.__self__._sldIdLst), "slides")


if __name__ == "__main__":
    main()
